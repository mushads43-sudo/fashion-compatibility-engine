from __future__ import annotations

from pathlib import Path
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def clean_markdown(text: str) -> list[str]:
    cleaned: list[str] = []
    in_code = False
    for raw_line in text.splitlines():
        line = raw_line.rstrip()
        if line.startswith("```"):
            in_code = not in_code
            continue
        line = line.replace("**", "").replace("`", "")
        line = re.sub(r"^\s*[-*]\s+", "- ", line)
        line = re.sub(r"^\s*\d+\.\s+", lambda match: match.group(0), line)
        if not line and cleaned and cleaned[-1] == "":
            continue
        cleaned.append(line)
    return cleaned


def split_lines(lines: list[str], max_chars: int = 1150, max_lines: int = 18) -> list[list[str]]:
    chunks: list[list[str]] = []
    current: list[str] = []
    current_chars = 0
    for line in lines:
        line_chars = len(line)
        if current and (current_chars + line_chars > max_chars or len(current) >= max_lines):
            chunks.append(current)
            current = []
            current_chars = 0
        current.append(line)
        current_chars += line_chars
    if current:
        chunks.append(current)
    return chunks


def add_title(slide, title: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(12.2), Inches(0.55))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title[:95]
    run.font.name = "Aptos"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(26, 34, 44)


def add_body(slide, lines: list[str]) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.95), Inches(12.0), Inches(5.95))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = line
        p.space_after = Pt(4)
        p.font.name = "Aptos"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(42, 48, 56)
        if line.startswith("- "):
            p.level = 0
            p.text = line
        elif line.endswith(":") or line.startswith("Title:") or line.startswith("Content:") or line.startswith("Speaker Notes:"):
            p.font.bold = True
            p.font.size = Pt(15)
            p.font.color.rgb = RGBColor(12, 111, 128)


def add_footer(slide, number: int) -> None:
    box = slide.shapes.add_textbox(Inches(11.8), Inches(7.05), Inches(1.0), Inches(0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(number)
    run.font.name = "Aptos"
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(120, 126, 136)


def add_slide(prs: Presentation, title: str, lines: list[str], number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title)
    add_body(slide, lines)
    add_footer(slide, number)


def section_title(section: str, fallback: str) -> str:
    first = section.splitlines()[0].strip()
    if first.startswith("Slide "):
        return first
    return fallback


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    outline_path = root / "docs" / "DETAILED_PPT_OUTLINE_FOR_SIR.md"
    output = root / "docs" / "complete_detailed_ppt_for_sir.pptx"

    text = outline_path.read_text(encoding="utf-8")
    parts = re.split(r"\n## ", text)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_number = 1

    intro = clean_markdown(parts[0])
    for chunk in split_lines(intro, max_chars=1000, max_lines=16):
        add_slide(prs, "Project Introduction", chunk, slide_number)
        slide_number += 1

    for part in parts[1:]:
        section = part.strip()
        if not section:
            continue
        title = section_title(section, "Project Details")
        lines = clean_markdown(section)
        if lines:
            lines = lines[1:]
        chunks = split_lines(lines, max_chars=1050, max_lines=17)
        for index, chunk in enumerate(chunks):
            chunk_title = title if index == 0 else f"{title} (continued)"
            add_slide(prs, chunk_title, chunk, slide_number)
            slide_number += 1

    prs.save(output)
    print(f"Saved complete PPT to {output}")


if __name__ == "__main__":
    main()
