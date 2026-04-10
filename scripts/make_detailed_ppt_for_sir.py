from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


TITLE = "Fashion Visual Similarity + Outfit Compatibility Engine"
SUBTITLE = "An AI system for visual product search and outfit compatibility scoring"


SLIDES = [
    {
        "title": TITLE,
        "subtitle": SUBTITLE,
        "bullets": ["Name", "Course/Subject", "Guide/Sir Name", "Department/College"],
        "layout": "title",
    },
    {
        "title": "Problem Statement",
        "bullets": [
            "Fashion platforms contain thousands of products.",
            "Users often cannot describe fashion style accurately using text.",
            "Traditional recommenders depend on clicks, ratings, or purchase history.",
            "They struggle with new users, new products, and visual style matching.",
            "This project finds similar products and predicts outfit compatibility using images.",
        ],
        "quote": "Given clothing images, the system finds visually similar products and predicts whether two fashion items aesthetically go well together.",
    },
    {
        "title": "Why This Project?",
        "bullets": [
            "Fashion is a highly visual domain.",
            "Users often shop by appearance, not only product name.",
            "Matching outfits requires color, category, pattern, and style understanding.",
            "Real companies use visual AI for search and complete-the-look systems.",
            "This project demonstrates a practical version of that system.",
        ],
        "quote": "blue checked shirt + navy jeans = good casual outfit",
    },
    {
        "title": "Project Objectives",
        "bullets": [
            "Use real-world fashion product images.",
            "Extract image features using a pretrained CNN.",
            "Find visually similar products using vector similarity.",
            "Group similar products into style clusters.",
            "Predict compatibility between two fashion items.",
            "Build a working Streamlit demo and FastAPI backend.",
        ],
    },
    {
        "title": "Main Features",
        "bullets": [
            "Visual Similarity Search: upload one clothing image and find visually similar products.",
            "Results include similarity scores and style clusters.",
            "Outfit Compatibility Prediction: upload two item images and get an aesthetic score.",
            "The output is shown as a percentage.",
        ],
        "quote": "Black polo t-shirt + blue shorts = 74.4% compatible",
    },
    {
        "title": "Real-World Datasets",
        "bullets": [
            "Fashion Product Images: Kaggle product photos and metadata, used for implementation.",
            "Polyvore Outfit Dataset: real outfit combinations for compatibility research.",
            "DeepFashion: large academic dataset for retrieval, attributes, and category prediction.",
            "Kaggle dataset is used as the main working dataset because it is practical and downloadable.",
        ],
    },
    {
        "title": "Fashion Product Images Dataset",
        "bullets": [
            "Dataset source: Kaggle Fashion Product Images.",
            "Files used: styles.csv and images/.",
            "Metadata: product ID, gender, category, article type, color, season, usage, product name.",
            "Contains real product images and useful metadata.",
            "Works well for visual search and local demo deployment.",
        ],
    },
    {
        "title": "Data Preparation",
        "bullets": [
            "Raw dataset: styles.csv + images/.",
            "Processed catalog: data/processed/catalog.csv.",
            "Clean fields: item ID, image path, gender, category, article type, color, season, usage, product name.",
            "Removes missing image entries.",
            "Creates clean input for the ML pipeline.",
        ],
    },
    {
        "title": "System Architecture",
        "bullets": [
            "Dataset -> catalog preparation.",
            "Catalog images -> ResNet50 image embeddings.",
            "Embeddings -> PCA compression.",
            "Compressed vectors -> KMeans style clustering.",
            "Vectors -> nearest-neighbor visual search.",
            "Pairwise embeddings -> compatibility classifier.",
            "Results -> Streamlit app and FastAPI backend.",
        ],
    },
    {
        "title": "Machine Learning Pipeline",
        "bullets": [
            "Load product images.",
            "Extract visual embeddings using ResNet50.",
            "Normalize embeddings.",
            "Compress embeddings using PCA.",
            "Cluster products using KMeans.",
            "Search similar products using cosine nearest neighbors.",
            "Train compatibility model using pairwise features.",
        ],
    },
    {
        "title": "Live Demo / Screenshots",
        "bullets": [],
        "blank": True,
    },
    {
        "title": "CNN Feature Extraction Using ResNet50",
        "bullets": [
            "CNNs are used for image understanding.",
            "ResNet50 is a pretrained deep CNN.",
            "The final classification layer is removed.",
            "The model outputs a 2048-dimensional image embedding.",
        ],
        "quote": "Fashion image -> ResNet50 -> 2048-dimensional vector",
    },
    {
        "title": "Image Embeddings",
        "bullets": [
            "An image embedding is a numerical vector representation of an image.",
            "Similar-looking images should have similar vectors.",
            "Embeddings capture color, texture, shape, pattern, product type, and style.",
            "Embeddings allow images to be compared mathematically instead of comparing raw pixels.",
        ],
        "quote": "shirt image -> [0.12, -0.04, 0.88, ..., 0.31]",
    },
    {
        "title": "PCA for Dimensionality Reduction",
        "bullets": [
            "ResNet50 embedding size is 2048 dimensions.",
            "PCA reduces embeddings to 128 or 256 dimensions.",
            "PCA keeps the most important visual information.",
            "Benefits: faster search, lower memory usage, reduced noise, easier clustering.",
        ],
        "quote": "2048-dimensional embedding -> PCA -> 128-dimensional embedding",
    },
    {
        "title": "Style Grouping Using KMeans",
        "bullets": [
            "KMeans is an unsupervised learning algorithm.",
            "It groups visually similar products into clusters.",
            "The model discovers groups automatically from embeddings.",
            "Possible clusters: sports shoes, casual shirts, jeans, dresses, women's tops.",
        ],
    },
    {
        "title": "How Similar Products Are Found",
        "bullets": [
            "User uploads a product image.",
            "ResNet50 extracts an embedding.",
            "PCA compresses the embedding.",
            "The system compares it with catalog embeddings.",
            "Nearest neighbors are found using cosine distance.",
            "Top matching products are returned.",
        ],
        "quote": "Query image -> embedding -> PCA vector -> cosine nearest neighbors -> top similar products",
    },
    {
        "title": "How Compatibility Is Predicted",
        "bullets": [
            "Compatibility is different from similarity.",
            "Two similar products may not form an outfit.",
            "The system compares two item embeddings using pairwise features.",
            "Features: cosine similarity, absolute difference, element-wise product.",
            "Classifier output: P(compatible | item A, item B).",
        ],
    },
    {
        "title": "Result Interpretation",
        "bullets": [
            "The compatibility model outputs a probability.",
            "0.744 means 74.4% compatible.",
            "70% and above: strong match.",
            "50% to 70%: good match.",
            "30% to 50%: possible match.",
            "Below 30%: weak match.",
        ],
    },
    {
        "title": "User Interface and API",
        "bullets": [
            "Streamlit app: upload images, view similar products, check outfit compatibility.",
            "FastAPI backend supports deployment and integration.",
            "API endpoints: /health, /similar, /compatibility.",
            "The project is usable as an interactive application, not only ML code.",
        ],
    },
    {
        "title": "Beyond Traditional Recommendation",
        "bullets": [
            "Traditional systems use clicks, ratings, purchase history, and co-purchase data.",
            "This project uses image content, visual embeddings, style clusters, and pairwise compatibility.",
            "Works for new products.",
            "Does not need user history.",
            "Understands product appearance and supports visual search.",
        ],
    },
    {
        "title": "Sample High-Compatibility Results",
        "bullets": [
            "17240.jpg + 54924.jpg: black striped polo + blue shorts, approx. 74.4%.",
            "49653.jpg + 51499.jpg: green women's top + blue jeans, approx. 66.9%.",
            "6617.jpg + 54924.jpg: black-white check shirt + blue shorts, approx. 59.2%.",
            "These examples can be used for live upload or screenshots.",
        ],
    },
    {
        "title": "Advantages of the Project",
        "bullets": [
            "Uses real-world fashion product images.",
            "Combines computer vision and machine learning.",
            "Provides both retrieval and compatibility scoring.",
            "Includes unsupervised style clustering.",
            "Has a Streamlit demo and FastAPI backend.",
            "Uses a clean GitHub-ready project structure.",
        ],
    },
    {
        "title": "Current Limitations",
        "bullets": [
            "Compatibility model uses metadata-based pair generation.",
            "Full Polyvore image training is limited because many original image URLs are unavailable.",
            "ResNet50 is general-purpose, not fashion-specific.",
            "Scores are useful for demo but not production-perfect.",
            "No user personalization is included yet.",
        ],
    },
    {
        "title": "Future Enhancements",
        "bullets": [
            "Use FashionCLIP or CLIP for better fashion embeddings.",
            "Train with complete human-curated outfit datasets.",
            "Add triplet loss or contrastive learning.",
            "Use FAISS or vector databases for large-scale search.",
            "Add occasion-based styling and color harmony.",
            "Build a complete outfit generation system.",
        ],
    },
    {
        "title": "Applications",
        "bullets": [
            "Fashion e-commerce product discovery.",
            "Visual search engines.",
            "Online stylist assistants.",
            "Complete-the-look recommendation.",
            "Product similarity detection.",
            "Catalog organization.",
            "Outfit planning apps.",
        ],
    },
    {
        "title": "Conclusion",
        "bullets": [
            "The project uses real fashion images for visual AI.",
            "ResNet50 embeddings represent product appearance.",
            "PCA and nearest-neighbor search find similar products.",
            "KMeans discovers visual style groups.",
            "Pairwise classification predicts outfit compatibility.",
            "Streamlit and FastAPI make the project usable and deployable.",
        ],
        "quote": "This project teaches a machine to understand fashion visually, so it can find similar products and predict which items aesthetically work together.",
    },
]


def set_run_style(run, size: int = 22, bold: bool = False, color: RGBColor | None = None) -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Aptos"
    if color:
        run.font.color.rgb = color


def add_title(slide, title: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.25), Inches(0.7))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    set_run_style(run, size=30, bold=True, color=RGBColor(28, 36, 46))


def add_bullets(slide, bullets: list[str], top: float = 1.35) -> None:
    box = slide.shapes.add_textbox(Inches(0.85), Inches(top), Inches(11.7), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for index, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(8)
        p.font.size = Pt(20 if len(bullets) <= 5 else 18)
        p.font.name = "Aptos"
        p.font.color.rgb = RGBColor(42, 48, 56)


def add_quote(slide, quote: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.85), Inches(6.15), Inches(11.7), Inches(0.8))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = quote
    set_run_style(run, size=18, bold=True, color=RGBColor(12, 111, 128))


def add_footer(slide, number: int) -> None:
    box = slide.shapes.add_textbox(Inches(11.95), Inches(7.05), Inches(0.8), Inches(0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(number)
    set_run_style(run, size=10, color=RGBColor(120, 126, 136))


def add_title_slide(prs: Presentation, slide_data: dict, number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, slide_data["title"])
    sub = slide.shapes.add_textbox(Inches(0.6), Inches(1.25), Inches(12), Inches(0.8))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = slide_data["subtitle"]
    set_run_style(run, size=24, color=RGBColor(12, 111, 128))
    add_bullets(slide, slide_data["bullets"], top=2.35)
    add_quote(slide, "We are teaching machines fashion sense, not just recommendations.")
    add_footer(slide, number)


def add_content_slide(prs: Presentation, slide_data: dict, number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, slide_data["title"])
    if slide_data.get("blank"):
        box = slide.shapes.add_textbox(Inches(1.2), Inches(2.7), Inches(10.9), Inches(1.2))
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "Leave this slide empty for live demo or screenshots"
        set_run_style(run, size=26, bold=True, color=RGBColor(145, 145, 145))
    else:
        add_bullets(slide, slide_data["bullets"])
        if slide_data.get("quote"):
            add_quote(slide, slide_data["quote"])
    add_footer(slide, number)


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    output = root / "docs" / "detailed_ppt_for_sir.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for index, slide_data in enumerate(SLIDES, start=1):
        if slide_data.get("layout") == "title":
            add_title_slide(prs, slide_data, index)
        else:
            add_content_slide(prs, slide_data, index)

    prs.save(output)
    print(f"Saved detailed PPT to {output}")


if __name__ == "__main__":
    main()
