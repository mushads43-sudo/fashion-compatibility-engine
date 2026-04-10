from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


SLIDES = [
    (
        "Fashion Visual Similarity + Outfit Compatibility Engine",
        ["We are teaching machines fashion sense, not just recommendations."],
    ),
    (
        "Problem",
        [
            "Fashion search is still too dependent on keywords and purchase history.",
            "New products suffer from cold-start recommendation problems.",
            "Outfit matching requires aesthetic understanding, not only user behavior.",
        ],
    ),
    (
        "Solution",
        [
            "Upload clothing images.",
            "Find visually similar catalog products.",
            "Predict whether two items go well together.",
            "Return a compatibility percentage for demo-ready styling intelligence.",
        ],
    ),
    (
        "Real Datasets",
        [
            "Fashion Product Images: real product photos and metadata.",
            "Polyvore Outfit Dataset: real outfit compatibility labels.",
            "DeepFashion: academic retrieval and attribute benchmark.",
        ],
    ),
    (
        "Architecture",
        [
            "Image upload -> ResNet50 embeddings -> PCA compression.",
            "PCA vectors -> KMeans style clusters -> nearest-neighbor search.",
            "Pair embeddings -> compatibility classifier -> score.",
        ],
    ),
    (
        "ML Concepts",
        [
            "CNN feature extraction with ResNet50.",
            "Cosine similarity for visual retrieval.",
            "PCA for efficient embeddings.",
            "KMeans for style grouping.",
            "Pairwise classification for compatibility.",
        ],
    ),
    (
        "Compatibility Model",
        [
            "Input: two clothing item embeddings.",
            "Features: cosine similarity, absolute difference, element-wise interaction.",
            "Output: P(compatible | item A, item B).",
        ],
    ),
    (
        "Deployment",
        [
            "FastAPI backend for inference.",
            "Streamlit app for image upload demo.",
            "Dockerfile for cloud deployment.",
            "Ready for Render, Railway, or Hugging Face Spaces.",
        ],
    ),
    (
        "Business Impact",
        [
            "Visual product discovery.",
            "Complete-the-look recommendations.",
            "Outfit scoring for styling assistants.",
            "Catalog deduplication and style intelligence.",
        ],
    ),
    (
        "Future Scope",
        [
            "FashionCLIP embeddings.",
            "Triplet-loss metric learning.",
            "Direct Polyvore image training.",
            "Vector database for millions of products.",
        ],
    ),
]


def add_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        paragraph = body.paragraphs[0] if i == 0 else body.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(24 if len(bullets) <= 2 else 20)


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    output = root / "docs" / "fashion_compatibility_engine.pptx"
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = SLIDES[0][0]
    title_slide.placeholders[1].text = SLIDES[0][1][0]

    for title, bullets in SLIDES[1:]:
        add_slide(prs, title, bullets)

    prs.save(output)
    print(f"Saved PPT to {output}")


if __name__ == "__main__":
    main()
